#!/usr/bin/env python3
"""Approximate PPTX -> PNG preview (layout check only, not a faithful renderer).
Draws shape rectangles/ovals with their fill+outline and lays out each run's text
with the same wrap rules, so overlaps and text overflow are visible."""
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU = 914400
DPI = 110
FONT = "/tmp/claude-1001/-workspace-speculators/97150579-9d65-4add-b551-73e1c59c966a/scratchpad/cjk.otf"
_cache = {}


def fnt(pt, bold=False):
    k = (round(pt), bold)
    if k not in _cache:
        _cache[k] = ImageFont.truetype(FONT, max(6, int(round(pt * DPI / 72))))
    return _cache[k]


def rgb(c, d=None):
    try:
        return tuple(bytes.fromhex(str(c)))
    except Exception:
        return d


def wrap(dr, txt, f, maxw):
    out, cur = [], ""
    for ch in txt:
        if dr.textlength(cur + ch, font=f) <= maxw or not cur:
            cur += ch
        else:
            out.append(cur); cur = ch
    if cur:
        out.append(cur)
    return out


def render(pptx, out_prefix, only=None):
    prs = Presentation(pptx)
    W = int(prs.slide_width / EMU * DPI); H = int(prs.slide_height / EMU * DPI)
    for idx, slide in enumerate(prs.slides, 1):
        if only and idx not in only:
            continue
        img = Image.new("RGB", (W, H), "white"); dr = ImageDraw.Draw(img)
        for sh in slide.shapes:
            x = int((sh.left or 0) / EMU * DPI); y = int((sh.top or 0) / EMU * DPI)
            w = int((sh.width or 0) / EMU * DPI); h = int((sh.height or 0) / EMU * DPI)
            fill = line = None
            try:
                if sh.fill.type is not None and sh.fill.type == 1:
                    fill = rgb(sh.fill.fore_color.rgb)
            except Exception:
                pass
            try:
                if sh.line.fill.type == 1:
                    line = rgb(sh.line.color.rgb)
            except Exception:
                pass
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or sh.has_table:
                if getattr(sh, "auto_shape_type", None) is not None and "OVAL" in str(sh.auto_shape_type):
                    dr.ellipse([x, y, x + w, y + h], fill=fill, outline=line, width=2)
                elif fill or line:
                    dr.rectangle([x, y, x + w, y + h], fill=fill, outline=line, width=2)
            if sh.has_table:
                t = sh.table
                cy = y
                for r in t.rows:
                    cx = x; rh = int(r.height / EMU * DPI)
                    for c in t.columns:
                        cw = int(c.width / EMU * DPI)
                        dr.rectangle([cx, cy, cx + cw, cy + rh], outline=(200, 205, 215))
                        cx += cw
                    cy += rh
                for ri, r in enumerate(t.rows):
                    pass
                cy = y
                for ri, r in enumerate(t.rows):
                    cx = x; rh = int(r.height / EMU * DPI)
                    for ci, c in enumerate(t.columns):
                        cw = int(c.width / EMU * DPI)
                        cell = t.cell(ri, ci); ty = cy + 3
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                f = fnt(run.font.size.pt if run.font.size else 11, run.font.bold)
                                for ln in wrap(dr, run.text, f, cw - 10):
                                    dr.text((cx + 5, ty), ln, font=f, fill=rgb(run.font.color.rgb, (0, 0, 0)) if run.font.color and run.font.color.type is not None else (0, 0, 0))
                                    ty += f.size + 2
                        cx += cw
                    cy += rh
                continue
            if sh.has_text_frame:
                ty = y + 3
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        f = fnt(run.font.size.pt if run.font.size else 12, run.font.bold)
                        col = (0, 0, 0)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                col = rgb(run.font.color.rgb, (0, 0, 0))
                        except Exception:
                            pass
                        for ln in wrap(dr, run.text, f, max(20, w - 10)):
                            dr.text((x + 5, ty), ln, font=f, fill=col)
                            ty += f.size + 2
                if ty > y + h + 6:
                    dr.rectangle([x, y, x + w, y + h], outline=(255, 0, 0), width=2)
        p = f"{out_prefix}_{idx}.png"; img.save(p); print("saved", p)


if __name__ == "__main__":
    only = {int(a) for a in sys.argv[3:]} if len(sys.argv) > 3 else None
    render(sys.argv[1], sys.argv[2], only)


def overlaps(pptx):
    """Flag shape pairs that overlap enough to look like a layout bug."""
    prs = Presentation(pptx)
    for idx, sl in enumerate(prs.slides, 1):
        bs = [(sh.shape_id, sh.left or 0, sh.top or 0, (sh.left or 0) + (sh.width or 0),
               (sh.top or 0) + (sh.height or 0)) for sh in sl.shapes]
        for i in range(len(bs)):
            for j in range(i + 1, len(bs)):
                a, b = bs[i], bs[j]
                ox = min(a[3], b[3]) - max(a[1], b[1])
                oy = min(a[4], b[4]) - max(a[2], b[2])
                # skip legitimate nesting: one shape fully inside the other
                inside = (a[1] <= b[1] and a[2] <= b[2] and a[3] >= b[3] and a[4] >= b[4]) or \
                         (b[1] <= a[1] and b[2] <= a[2] and b[3] >= a[3] and b[4] >= a[4])
                if inside:
                    continue
                if ox > 0 and oy > EMU * 0.06:          # >0.06in vertical bite
                    amin = min((a[3]-a[1])*(a[4]-a[2]), (b[3]-b[1])*(b[4]-b[2]))
                    if ox * oy > 0.25 * amin:
                        print(f"  slide {idx}: shapes {a[0]}&{b[0]} overlap "
                              f"{ox/EMU:.2f}x{oy/EMU:.2f}in")


def text_overflow(pptx):
    """Report text frames whose laid-out text runs past their own shape."""
    from PIL import Image, ImageDraw
    prs = Presentation(pptx)
    dr = ImageDraw.Draw(Image.new("RGB", (10, 10)))
    for idx, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_text_frame or sh.has_table:
                continue
            w = int((sh.width or 0) / EMU * DPI); h = int((sh.height or 0) / EMU * DPI)
            used = 6
            for p in sh.text_frame.paragraphs:
                for run in p.runs:
                    f = fnt(run.font.size.pt if run.font.size else 12, run.font.bold)
                    used += len(wrap(dr, run.text, f, max(20, w - 10))) * (f.size + 2)
            if used > h + 6:
                print(f"  slide {idx}: shape {sh.shape_id} text {used/DPI:.2f}in "
                      f"> box {h/DPI:.2f}in — {sh.text_frame.text[:34]!r}")
