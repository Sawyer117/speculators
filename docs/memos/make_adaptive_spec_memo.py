#!/usr/bin/env python3
"""生成技术备忘录 PPT:自适应投机解码在图模式下的代价(昇腾 NPU)。

结论均有代码/文档出处:vllm-ascend 源码、vLLM PR #48692、昇腾算子文档。
用法:python3 make_adaptive_spec_memo.py --out <file>.pptx
"""
from __future__ import annotations

import argparse
import copy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1B, 0x22, 0x2E)
MUTE = RGBColor(0x5B, 0x66, 0x77)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x0E, 0x7A, 0x4B)
WARN = RGBColor(0xB4, 0x53, 0x09)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF2, 0xF5, 0xF9)
CJK = "微软雅黑"


def _cjk(run, name=CJK):
    """python-pptx 只设 a:latin;中文需要同时设 a:ea,否则 PowerPoint 会回退字体。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = copy.deepcopy(latin) if latin is not None else None
        if ea is None:
            return
        ea.tag = qn("a:ea")
        rPr.append(ea)
    ea.set("typeface", name)


def _txbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    return tf


def _p(tf, text, *, size=15, bold=False, color=INK, before=6, after=0, level=0,
       align=PP_ALIGN.LEFT, first=False, mono=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment, para.level = align, level
    para.space_before, para.space_after = Pt(before), Pt(after)
    run = para.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    _cjk(run, "Consolas" if mono else CJK)
    return para


def _slide(prs, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.95))
    _p(tf, title, size=25, bold=True, first=True, before=0)
    if kicker:
        _p(tf, kicker, size=12, color=MUTE, before=3)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.22), Inches(1.3), Emu(24000))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background(); ln.shadow.inherit = False
    return s


def _table(slide, rows, x, y, w, col_w=None, font=11.5, row_h=0.32, colors=None):
    t = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, Inches(row_h * len(rows))).table
    if col_w:
        tot = sum(col_w)
        for i, f in enumerate(col_w):
            t.columns[i].width = Emu(int(w * f / tot))
    t.first_row = True
    for r, row in enumerate(rows):
        t.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if r == 0 else BG
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.text_frame.word_wrap = True
            for i, line in enumerate(str(val).split("\n")):
                para = cell.text_frame.paragraphs[0] if i == 0 else cell.text_frame.add_paragraph()
                para.space_before = para.space_after = Pt(0)
                run = para.add_run()
                run.text = line
                run.font.size = Pt(font)
                run.font.bold = (r == 0)
                col = INK
                if colors and r > 0 and (r, c) in colors:
                    col = colors[(r, c)]
                run.font.color.rgb = col
                _cjk(run)
    return t


def _code(slide, lines, x, y, w, h, font=11):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = BAND
    box.line.color.rgb = RGBColor(0xD8, 0xDF, 0xE8); box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left, tf.margin_top = Inches(0.14), Inches(0.08)
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = para.space_after = Pt(0)
        run = para.add_run()
        run.text = ln
        run.font.size = Pt(font)
        run.font.color.rgb = INK
        _cjk(run, "Consolas")
    return box


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 封面 ────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.8))
    _p(tf, "自适应投机解码的图模式代价", size=38, bold=True, first=True, before=0)
    _p(tf, "昇腾 NPU 上的真实瓶颈在哪 —— 通用模型与 DSV4 两条路线", size=19, color=MUTE, before=10)
    _p(tf, "技术备忘录 · 结论出自 vllm-ascend 源码、vLLM PR #48692、昇腾算子文档", size=12, color=MUTE, before=20)

    # 2 结论速览 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "结论速览", "自适应投机按置信度给每个请求分配 K_i 以少算 target token;设备图要求形状/地址/执行路径固定。问题不是“能不能省”,而是“为适配图付出的额外开销是否吃掉收益”")
    _table(s, [
        ["", "结论", "依据"],
        ["① GPU 参考被高估",
         "PR #48692 已关闭未合入;明确不做在线预算分配;仅 FLASH_ATTENTION;\n作者自述测不出加速,收益体现在接受率",
         "PR 正文"],
        ["② 设备侧长度已存在\n(但只在一族算子)",
         "DSA / lightning-indexer / SFA 收 torch.Tensor;FIA v2 收 host int 数组。\nDSV4 走 DSA → 前置条件已满足;通用稠密模型走 FIA → 卡住",
         "device_op.py 签名\n昇腾算子文档"],
        ["③ 增量代价被误判",
         "graph_task_update 固定-K 投机已在付;动态 K 的真正增量是\nworkspace 每轮 miss + 图键爆炸 + padding",
         "attention_v1.py\n图 replay 路径"],
    ], Inches(0.55), Inches(1.75), Inches(12.3), col_w=[0.17, 0.60, 0.23], row_h=0.80,
       colors={(2, 1): GOOD, (3, 1): WARN})
    tf = _txbox(s, Inches(0.55), Inches(5.35), Inches(12.3), Inches(1.6))
    _p(tf, "建议:先量化“动态分配到底少算多少 target token”,再决定是否啃图路径。",
       size=17, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "GPU 侧原型在 8B 上做,verify 本来就便宜、压根没证明收益;收益随 target 变贵而放大 —— 大 MoE target 才是值得做的场景。",
       size=13.5, color=MUTE, before=8)

    # 3 PR #48692 ───────────────────────────────────────────────────────────
    s = _slide(prs, "① GPU 参考:PR #48692 到底交付了什么", "vLLM main · [MRV2][Spec Decode] Adaptive Speculative Decoding · closed / needs-rebase · 28 文件 +1240/−110")
    _table(s, [
        ["", "通常的理解", "PR 实际写的"],
        ["验证预算", "confidence head 在线分配", "明确不实现;改用用户配置的\nnum_speculative_tokens_per_batch_size"],
        ["加速", "已验证", "“no significant speedup is measurable;\nthat is not the objective”"],
        ["后端", "通用", "仅 FLASH_ATTENTION 可用"],
    ], Inches(0.55), Inches(1.72), Inches(12.3), col_w=[0.15, 0.30, 0.55], row_h=0.60)
    tf = _txbox(s, Inches(0.55), Inches(3.75), Inches(12.3), Inches(1.2))
    _p(tf, "它真正交付的:变长 per-request 投机与 FULL CUDA Graph 共存 —— attention metadata 的变长兼容标志、"
           "structured outputs、logprobs、prefill/decode 混合 batch。", size=14.5, first=True, before=0)
    _table(s, [
        ["SPEED-Bench 平均接受长度", "静态 K=5", "静态 K=7", "自适应(上限7/有效5)"],
        ["总体", "3.445", "3.918", "3.707"],
    ], Inches(0.55), Inches(5.05), Inches(8.6), col_w=[0.34, 0.22, 0.22, 0.22], row_h=0.34)
    tf = _txbox(s, Inches(0.55), Inches(5.95), Inches(12.3), Inches(1.0))
    _p(tf, "⟹ 用 K=5 的预算拿到 K=5→K=7 之间约 55% 的接受率增益。这是它唯一被证实的收益,与延迟无关。",
       size=14.5, bold=True, color=ACCENT, first=True, before=0)

    # 4 两条算子路径 ─────────────────────────────────────────────────────────
    s = _slide(prs, "② 两条算子路径的分野 —— 决定通用模型与 DSV4 难度不同", "“K_i 能否全程留在设备上”取决于走哪族算子,而不取决于硬件")
    _table(s, [
        ["", "FIA v2\nnpu_fused_infer_attention_score_v2", "DSA / lightning-indexer / SFA\ntorch.ops._C_ascend.npu_vllm_*"],
        ["长度契约", "Host 侧 int 数组(SymInt[])", "torch.Tensor(设备侧)"],
        ["代码证据", "全部调用点 .tolist() 构造\nactual_seq_lengths_q: list[int]", "execute_sparse_flash_attention_process(\n  actual_seq_lengths_query: torch.Tensor, …)\n调用处传 query_start_loc[1:].clone()"],
        ["适用模型", "通用稠密模型(Qwen3 等)", "DeepSeek-V4-Flash"],
        ["动态 K 前置条件", "不满足 —— 每轮长度必过 host", "已满足 —— 无需新造"],
    ], Inches(0.55), Inches(1.72), Inches(12.3), col_w=[0.16, 0.40, 0.44], row_h=0.72,
       colors={(4, 1): WARN, (4, 2): GOOD})
    tf = _txbox(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(1.5))
    _p(tf, "对照:GPU 侧 #48692 同样只覆盖 FLASH_ATTENTION —— 两边都是“先在一族算子上跑通,再谈通用”。",
       size=14.5, first=True, before=0)
    _p(tf, "⟹ DSV4 可以现在就做;通用模型这条线的前置工作是让 FIA v2 接受设备张量长度(算子层改动),或接受每轮 host metadata 的代价。",
       size=15, bold=True, color=ACCENT, before=8)

    # 5 增量代价 ─────────────────────────────────────────────────────────────
    s = _slide(prs, "③ 动态 K 的增量代价在哪 —— 不是 graph update", "固定-K 投机早已入图,replay 前逐层更新 seq lengths 的开销已经在付了")
    _code(s, [
        "# attention_v1.py —— 固定-K 投机的图 replay 已存在",
        "if _EXTRA_CTX.is_draft_model:",
        "    draft_step, key = draft_attn_key_steps[attn_count]",
        "    actual_seq_lengths_q = attn_metadata[draft_step][key].actual_seq_lengths_q",
        "torch.npu.graph_task_update_begin(update_stream, handle)   # 逐层更新 → 已有开销",
        "",
        "# 但 workspace 缓存键正好是动态 K 会变的那个量",
        "num_tokens = attn_metadata.actual_seq_lengths_q[-1]    # 本轮总 token 数",
        "workspace  = graph_params.workspaces.get(num_tokens)   # 按总 token 数缓存",
        "elif workspace is None:                                # 仅未命中时查询",
        "    workspace = torch_npu._npu_fused_infer_attention_score_v2_get_max_workspace(...)",
    ], Inches(0.55), Inches(1.72), Inches(12.3), Inches(2.55))
    _table(s, [
        ["增量项", "固定 K", "动态 K_i", "说明"],
        ["workspace 查询", "恒命中", "每轮每层 miss", "缓存键 = 总 token 数"],
        ["graph task update", "已在付", "同左,无增量", "两者都要逐层更新"],
        ["图键 / 捕图数量", "1 个形状", "按 K_i 组合爆炸", "必须改为按 bucket 建键"],
        ["padding 浪费", "无", "取决于 bucket 宽度", "换取图可复用的代价"],
    ], Inches(0.55), Inches(4.45), Inches(12.3), col_w=[0.22, 0.15, 0.25, 0.38], row_h=0.40,
       colors={(1, 2): WARN, (2, 2): GOOD, (3, 2): WARN})
    tf = _txbox(s, Inches(0.55), Inches(6.35), Inches(12.3), Inches(0.9))
    _p(tf, "修法是结构性的:缓存键与图键都改为 (batch bucket × 总验证 token bucket),捕图阶段一次性预分配 workspace。",
       size=15, bold=True, color=ACCENT, first=True, before=0)

    # 6 两条路线的适配清单 ────────────────────────────────────────────────────
    s = _slide(prs, "适配清单 —— 两条路线分开推进", "共用项做一次,路线项各自解决")
    _table(s, [
        ["", "共用(先做)", "DSV4 / DSA 路线", "通用 / FIA 路线"],
        ["工作项",
         "· workspace 与图键改为 bucket\n· 捕图期预分配并长期复用\n· K_i→有效 token 的 compaction 全留设备\n· 审计 .cpu() / .tolist() / 隐式 D2H",
         "· 直接复用现有设备张量长度\n· 验证 indexer/SAS 在变长下的\n  metadata 复用\n· 端到端跑通并计时",
         "· FIA v2 接受设备张量长度\n  (算子层改动,或等上游)\n· 否则按 bucket 摊薄 host metadata\n· 对齐 #48692 的变长兼容标志"],
        ["难度", "中", "低 —— 前置条件已满足", "高 —— 依赖算子接口"],
    ], Inches(0.55), Inches(1.72), Inches(12.3), col_w=[0.10, 0.32, 0.29, 0.29], row_h=1.30,
       colors={(2, 2): GOOD, (2, 3): WARN})
    tf = _txbox(s, Inches(0.55), Inches(5.5), Inches(12.3), Inches(1.5))
    _p(tf, "顺序建议:先在 DSV4/DSA 上把动态 K 端到端跑通并量化收益 —— 它前置条件最全、能最快回答“值不值得”;"
           "结论为正再推动 FIA v2 的接口改动,让通用模型受益。", size=15, first=True, before=0)

    # 7 待测清单 ─────────────────────────────────────────────────────────────
    s = _slide(prs, "动手前先测这四件事", "任何一项为负,后面的工程都不必做")
    _table(s, [
        ["", "测什么", "判断标准"],
        ["1 · 先量化收益", "同一流量下,静态 K 与自适应分配各自的 target 实际验证 token 数", "减少量小 → 图工程再省也不回本,直接停"],
        ["2 · 四种模式对比", "静态K全图 / 动态K eager / 动态K piecewise / 动态K全图。\n统计:验证 token 数、图命中率、padding 数、host metadata 时间、\nworkspace+tiling 时间、graph update 时间、attention 时间、端到端时延", "端到端时延必须优于静态K全图,\n否则收益被图外开销吃掉"],
        ["3 · 审计 host 往返", "confidence head → K_i → prefix 选择 → index compaction →\nposition / slot mapping → query length,全链路是否留在设备", "出现 .cpu() / .tolist() / 隐式 D2H\n即为阻断点"],
        ["4 · 确定分桶策略", "按 (batch bucket, 总验证 token bucket) 建图键与 workspace 键", "捕图数量可控 且 padding 浪费可接受"],
    ], Inches(0.55), Inches(1.72), Inches(12.3), col_w=[0.16, 0.51, 0.33], row_h=1.02)
    tf = _txbox(s, Inches(0.55), Inches(6.15), Inches(12.3), Inches(1.0))
    _p(tf, "一句话:自适应投机在算法上一定省 token,但省下的是否大于为固定形状付出的代价 —— 这是个必须用数据回答的问题,不是设计问题。",
       size=15, bold=True, color=ACCENT, first=True, before=0)

    prs.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="adaptive-spec-graph-mode-ascend.pptx")
    build(ap.parse_args().out)
