#!/usr/bin/env python3
"""生成 v3 立项建议书 PPT:昇腾 NPU 自适应投机解码。

v1 = 技术备忘录(调查结论);v2 = 首版立项书;v3 = 依据上游 vLLM #47808 / #48692 的
实际进展重写 —— 算法已由上游实现,项目重定义为昇腾侧适配。三版并存,不覆盖。
用法:python3 make_adaptive_spec_proposal_v3.py --out <file>.pptx
"""
from __future__ import annotations

import argparse
import copy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1B, 0x22, 0x2E)
MUTE = RGBColor(0x5B, 0x66, 0x77)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x0E, 0x7A, 0x4B)
WARN = RGBColor(0xB4, 0x53, 0x09)
RISK = RGBColor(0xA3, 0x1D, 0x2A)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF2, 0xF5, 0xF9)
CJK = "微软雅黑"


def _cjk(run, name=CJK):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        if latin is None:
            return
        ea = copy.deepcopy(latin)
        ea.tag = qn("a:ea")
        rPr.append(ea)
    ea.set("typeface", name)


def _txbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    return tf


def _p(tf, text, *, size=15, bold=False, color=INK, before=6, after=0,
       align=PP_ALIGN.LEFT, first=False, mono=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before, para.space_after = Pt(before), Pt(after)
    run = para.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    _cjk(run, "Consolas" if mono else CJK)
    return para


def _slide(prs, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.55), Inches(0.28), Inches(12.3), Inches(0.95))
    _p(tf, title, size=25, bold=True, first=True, before=0)
    if kicker:
        _p(tf, kicker, size=12, color=MUTE, before=3)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.20), Inches(1.3), Emu(24000))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background(); ln.shadow.inherit = False
    return s


def _table(slide, rows, x, y, w, col_w=None, font=11.5, row_h=0.32, colors=None, bolds=None):
    t = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, Inches(row_h * len(rows))).table
    if col_w:
        tot = sum(col_w)
        for i, f in enumerate(col_w):
            t.columns[i].width = Emu(int(w * f / tot))
    t.first_row = True
    for r, row in enumerate(rows):
        t.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if r == 0 else BG
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.text_frame.word_wrap = True
            for i, line in enumerate(str(val).split("\n")):
                para = cell.text_frame.paragraphs[0] if i == 0 else cell.text_frame.add_paragraph()
                para.space_before = para.space_after = Pt(0)
                run = para.add_run()
                run.text = line
                run.font.size = Pt(font)
                run.font.bold = (r == 0) or bool(bolds and (r, c) in bolds)
                run.font.color.rgb = colors[(r, c)] if (colors and (r, c) in colors) else INK
                _cjk(run)
    return t




# ── 「技术点一页纸」总览页(参照华为内部单页汇报版式) ──────────────────────
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

RED = RGBColor(0xC0, 0x00, 0x00)
NAVY = RGBColor(0x1F, 0x30, 0x50)
BLUE = RGBColor(0x1F, 0x5F, 0xA8)
ORNG = RGBColor(0xD3, 0x62, 0x0B)
GREY = RGBColor(0xF2, 0xF2, 0xF2)
LGREY = RGBColor(0xE9, 0xEE, 0xF5)


def _shape(s, kind, x, y, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def _fit(sh, lines, size=9, color=INK, bold_first=False, align=PP_ALIGN.LEFT, space=1):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = p.space_after = Pt(space)
        txt, bold, col, sz = (ln if isinstance(ln, tuple) else (ln, False, color, size))
        r = p.add_run(); r.text = txt
        r.font.size, r.font.bold, r.font.color.rgb = Pt(sz), bold or (i == 0 and bold_first), col
        _cjk(r)
    return sh


def _tbox(s, x, y, w, h, lines, size=9, color=INK, space=1.5):
    tf = _txbox(s, Inches(x), Inches(y), Inches(w), Inches(h))
    for i, ln in enumerate(lines):
        txt, bold, col, sz = (ln if isinstance(ln, tuple) else (ln, False, color, size))
        _p(tf, txt, size=sz, bold=bold, color=col, first=(i == 0), before=(0 if i == 0 else space))
    return tf


def one_pager(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

    # ── 标题 ───────────────────────────────────────────────────────────────
    _tbox(s, 0.32, 0.14, 12.7, 0.5,
          [("技术点:昇腾 NPU 自适应投机解码 —— 按置信度分配验证预算", True, RED, 23)])

    # ── 挑战 / 目标 ────────────────────────────────────────────────────────
    Y, H = 0.70, 0.98
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, 0.32, Y, 0.42, H, RED),
         [("挑", False, BG, 12), ("战", False, BG, 12)], align=PP_ALIGN.CENTER, space=0)
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, 0.74, Y, 5.62, H, GREY),
         [("· 固定 K 对全 batch 一视同仁 —— 高置信请求本可多验,低置信请求白耗 target 算力", False, INK, 10.5),
          ("★ 变长 K_i 与设备图「形状固定」冲突 —— 若进不了图,即便自适应实现了,性能仍可能不如固定 K", False, RED, 10.5),
          ("· confidence head 在昇腾侧尚未加载(短期缺项,移植上游改法即可,非难点)", False, MUTE, 10.5)])
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, 6.62, Y, 0.42, H, RED),
         [("目", False, BG, 12), ("标", False, BG, 12)], align=PP_ALIGN.CENTER, space=0)
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, 7.04, Y, 5.94, H, GREY),
         [("· 对齐上游 vLLM #47808 机制,在昇腾实现按置信度的验证预算分配", False, INK, 10.5),
          ("· 同等验证预算下提升接受长度 —— 上游实测 +7.6%,收益随 target 变贵而放大", False, RED, 10.5),
          ("· 动态 K 在图模式下端到端不劣于同预算的固定 K", False, INK, 10.5)])

    # ── 左栏:运行闭环 ─────────────────────────────────────────────────────
    LX, LW = 0.32, 2.52
    _fit(_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, LX, 1.86, LW, 0.34, NAVY),
         [("自适应验证运行闭环", False, BG, 11.5)], align=PP_ALIGN.CENTER, space=0)
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, LX, 2.28, LW, 0.34, LGREY),
         [("输入:草稿 block + 逐位 confidence", False, INK, 8.5)], align=PP_ALIGN.CENTER, space=0)
    # 每个闭环把自己的「失败回退路径」写在框内,紧贴它所属的环 —— 不堆到列底下
    for (py, col, hdr, steps, fallback) in [
        (2.72, BLUE, "① 置信度调度闭环",
         ("生存概率打分", "全局预算择优", "每请求 K_i"), ""),
        (4.82, ORNG, "② 变长图执行闭环",
         ("变长 decode 图捕获", "真机成本曲线标定", "场景配置"),
         "⚠ 变长进不了图 → 退回 eager,失去图重放收益"),
    ]:
        _shape(s, MSO_SHAPE.RECTANGLE, LX, py, LW, 1.92, None, col, 1.25)
        _fit(_shape(s, MSO_SHAPE.RECTANGLE, LX + 0.06, py + 0.06, LW - 0.12, 0.28, None),
             [(hdr, True, col, 10)], space=0)
        for j, t in enumerate(steps):
            _fit(_shape(s, MSO_SHAPE.RECTANGLE, LX + 0.14, py + 0.38 + j * 0.42, LW - 0.28, 0.32,
                        BG, col, 0.75), [(t, False, INK, 9)], align=PP_ALIGN.CENTER, space=0)
            if j < 2:
                _shape(s, MSO_SHAPE.DOWN_ARROW, LX + LW / 2 - 0.05, py + 0.71 + j * 0.42,
                       0.10, 0.08, col)
        if fallback:
            _tbox(s, LX + 0.12, py + 1.60, LW - 0.24, 0.28, [(fallback, False, RED, 7.5)])
    _shape(s, MSO_SHAPE.DOWN_ARROW, LX + LW / 2 - 0.05, 4.68, 0.10, 0.10, INK)
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, LX, 6.86, LW, 0.34, LGREY),
         [("可部署:动态 K 投机解码", True, INK, 9.5)], align=PP_ALIGN.CENTER, space=0)

    # ── 中栏 ① 算法侧 ─────────────────────────────────────────────────────
    MX, MW = 2.98, 7.42
    _shape(s, MSO_SHAPE.RECTANGLE, MX, 1.86, MW, 2.58, None, BLUE, 1.5)
    _tbox(s, MX + 0.12, 1.94, 3.9, 0.34, [("① 置信度调度(算法侧,对标 #47808)", True, BLUE, 13)])
    _tbox(s, MX + 0.12, 2.34, 3.85, 1.62, [
        ("A. 生存概率打分", True, INK, 10),
        ("每个 (请求, 位置) 的 draft slot 以该请求逐位 confidence 的累积乘积打分。", False, INK, 9),
        ("B. 跨请求竞争的全局预算", True, INK, 10),
        ("所有 slot 统一排序、择优录取直至预算耗尽 —— 自信请求的第 5 位可压过犹豫请求的第 1 位。", False, INK, 9),
        ("C. 成本曲线标定", True, INK, 10),
        ("启动期 dummy step 测出步开销;batch 预算用 CPU 侧陈旧值(不同步),per-request 用 GPU 实时值。", False, INK, 9)], space=2)
    for j, (t, sub) in enumerate([("draft slot", "逐位置信度"), ("排序 / 打分", "累积乘积"), ("录取前缀", "预算内择优")]):
        bx = MX + 4.12 + j * 1.08
        _fit(_shape(s, MSO_SHAPE.RECTANGLE, bx, 2.62, 0.92, 0.86, BG, BLUE, 0.75),
             [(t, True, BLUE, 9), (sub, False, MUTE, 8)], align=PP_ALIGN.CENTER, space=1)
        if j < 2:
            _shape(s, MSO_SHAPE.RIGHT_ARROW, bx + 0.94, 2.97, 0.12, 0.14, BLUE)
    _tbox(s, MX + 4.12, 3.60, 3.16, 0.5, [("机制示意:同预算下把验证 token 挪给最可能被接受的位置", False, MUTE, 8.5)])
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, MX + 0.12, 4.02, MW - 0.24, 0.30, LGREY),
         [("交付:昇腾预算调度 kernel  |  成本曲线标定脚本  |  置信度-接受率相关性报告", True, BLUE, 9.5)],
         align=PP_ALIGN.CENTER, space=0)

    # ── 中栏 ② 工程侧 ─────────────────────────────────────────────────────
    _shape(s, MSO_SHAPE.RECTANGLE, MX, 4.60, MW, 2.72, None, ORNG, 1.5)
    _tbox(s, MX + 0.12, 4.68, 4.2, 0.34, [("② 变长 decode 图(工程侧,★ 主要投入)", True, ORNG, 13)])
    for j, (t, sub) in enumerate([("固定形状捕获", "现状"), ("变长图", "按 token 网格"), ("真机测量", "开销 / padding")]):
        bx = MX + 0.16 + j * 1.24
        _fit(_shape(s, MSO_SHAPE.RECTANGLE, bx, 5.10, 1.08, 0.80, BG, ORNG, 0.75),
             [(t, True, ORNG, 9), (sub, False, MUTE, 8)], align=PP_ALIGN.CENTER, space=1)
        if j < 2:
            _shape(s, MSO_SHAPE.RIGHT_ARROW, bx + 1.10, 5.43, 0.12, 0.14, ORNG)
    _tbox(s, MX + 0.16, 5.96, 3.74, 0.92, [
        ("让变长进图有两条路线,阶段 2 先验证可行性:", True, INK, 8),
        ("· 对齐上游 —— 图按 token 总数分档捕获,运行时用批内最长请求的长度挑图、其余补齐;"
         "图数 = 档数 × 长度取值数,有补齐浪费;", False, INK, 7),
        ("· SGLang 式 —— 各请求 token 紧挨着打包成一条,图只按总 token 数分档;"
         "图少且无补齐浪费,代价是改数据布局 + 算子须支持。", False, INK, 7)], space=1.5)
    _tbox(s, MX + 4.02, 4.94, 3.30, 2.10, [
        ("· 固定 K 每步形状一致 → 能入图;自适应下每请求 K_i 每步不同 → 形状参差且逐步变化 → 默认出图;",
         False, INK, 8),
        ("· 能否留在图内,取决于图里的长度参数是「录制时固化的常量」还是「运行时可改的设备张量」:",
         False, INK, 8),
        ("  通用路径 FIA v2(npu_fused_infer_attention_score_v2)—— 长度是 host 侧 list,"
         "录制即固化 → K 一变必须整图重录;", False, INK, 8),
        ("  ★ DSV4 路径 DSA / SFA —— 长度是 device tensor,图里存的是指针 → 改内容即可,"
         "结构上无需重录(有利条件);", False, INK, 8),
        ("· ★ 真正的未知量:workspace(算子临时显存)按 token 数做缓存键,K 一变即未命中、重新分配 →"
         " 图与 workspace 数量随 K 组合膨胀,显存与捕获耗时随之上升。", False, RED, 8)], space=2)
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, MX + 0.12, 6.92, MW - 0.24, 0.30, RGBColor(0xFD, 0xF0, 0xE4)),
         [("交付:变长 decode 图实现  |  Shape 配置  |  三模式(无投机 / 固定 K / 动态 K)对比实测包", True, ORNG, 9.5)],
         align=PP_ALIGN.CENTER, space=0)

    # ── 右栏:里程碑 ───────────────────────────────────────────────────────
    RX, RW = 10.54, 2.44
    _fit(_shape(s, MSO_SHAPE.RECTANGLE, RX, 1.86, RW, 0.34, RED),
         [("里 程 碑", True, BG, 12)], align=PP_ALIGN.CENTER, space=0)
    _shape(s, MSO_SHAPE.RECTANGLE, RX + 0.22, 2.42, 0.028, 4.55, RED)
    MS = [("T+1 天", "G1 · 判定", ["把加速比从 c=48 补测到", "c=128 / c=256", "看固定 K 是否收益衰减", "不达标即停,不进入开发"]),
          ("T+3 周", "G2 · 打通信号", ["vllm-ascend 加载置信度头", "端到端读出逐位置信度", "与实测接受率做相关性验证"]),
          ("T+11 周", "G3 · 图模式 ★", ["动态 K 在图模式下运行", "端到端不劣于同预算固定 K", "本项目主要技术风险所在"]),
          ("T+15 周", "合入与扩展", ["推广到通用模型路径", "补齐限制项 + 复现文档"])]
    for j, (when, what, det) in enumerate(MS):
        y = 2.46 + j * 1.16
        _fit(_shape(s, MSO_SHAPE.OVAL, RX + 0.09, y, 0.30, 0.30, BG, RED, 1.25),
             [(str(j + 1), True, RED, 10)], align=PP_ALIGN.CENTER, space=0)
        _tbox(s, RX + 0.46, y - 0.06, RW - 0.5, 1.1,
              [(when, True, RED, 11.5), (what, True, INK, 10)] + [(d, False, MUTE, 8) for d in det], space=1)
    _tbox(s, RX, 7.06, RW, 0.34, [("T = 立项批准日;绝对日期待排期确定后填入", False, MUTE, 7.5)])
    return s


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 封面 ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.2))
    _p(tf, "昇腾 NPU 自适应投机解码", size=38, bold=True, first=True, before=0)
    _p(tf, "立项建议书 v3", size=22, color=ACCENT, before=8)
    _p(tf, "按验证长度的置信度分配,替代全 batch 统一的固定 K", size=17, color=MUTE, before=12)
    _p(tf, "v3 相对 v2 的关键变化:上游 vLLM 已实现该算法(PR #47808,针对 DeepSeek-V4-Flash-DSpark),"
           "SGLang 亦已上线。本项目据此由“设计并实现”重定义为“昇腾侧适配与验证”,范围与周期同步收缩。",
       size=12, color=MUTE, before=20)

    # 2 技术点一页纸(总览) ───────────────────────────────────────────────────
    one_pager(prs)

    # 3 立项理由 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "一、为什么立项", "固定 K 对所有请求一视同仁,而“该验多长”的信号其实已经有了")
    _table(s, [
        ["", "现状", "问题", "本项目要解决的"],
        ["算法侧", "全 batch 统一验证长度 K", "高置信请求本可多验,低置信请求\n白白消耗 target 计算",
         "按 confidence head 的逐位接受概率\n给每请求分配 K_i"],
        ["工程侧", "固定 K 已入图,收益稳定", "K_i 每轮变化 → ragged batch,\n与设备图“形状固定”天然冲突",
         "变长与图共存,且开销可控"],
        ["生态侧", "GPU 侧上游已实现并在合入中\nNPU 侧空白", "上游方案绑定 NVIDIA 后端能力,\n昇腾无对应机制",
         "昇腾侧首个可用实现"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.09, 0.25, 0.33, 0.33], row_h=0.72)

    tf = _txbox(s, Inches(0.55), Inches(4.35), Inches(12.3), Inches(2.7))
    _p(tf, "核心依据:同等验证预算下,接受长度提升 7.6%(上游实测)", size=17, bold=True, color=ACCENT,
       first=True, before=0)
    _p(tf, "· vLLM PR #48692 在 Qwen3-8B + DSpark 上测得:固定 K=5 接受长度 3.445;自适应(上限 K=7、"
           "每请求有效 K=5,即同等预算)3.707 —— 同样的算力,多接受 7.6% 的 token。", size=13.5, before=8)
    _p(tf, "· 该 PR 作者同时说明:8B 模型上“测不出显著加速”,因为 verify 本来就便宜。", size=13.5, before=4)
    _p(tf, "  ⟹ 收益随 target 变贵而放大。target 越大、verify 越贵,少验一个 token 省下的越多 —— "
           "大 MoE(DSV4-Flash 284B/13B 激活)正是收益最大的场景,也正是我们已有完整链路的场景。",
       size=13.5, bold=True, before=4)
    _p(tf, "· 边界情形(非主因,但值得记录):并发继续升高时固定 K 会由正收益转为负收益。上游在 c=256 实测"
           "固定 7-token 比不开投机还慢 33%,自适应可保住收益。此时自适应从“更好”变为“必需”。",
       size=12.5, color=MUTE, before=8)

    # 3 上游现状 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "二、上游已经做到哪一步 —— 我们不再需要设计算法",
               "两个 PR 是同一条线的两代;算法、调度、图支持均已成型")
    _table(s, [
        ["", "vLLM #48692", "vLLM #47808(当前主线)"],
        ["状态", "2026-08-05 关闭\n作者:改推 #47808", "OPEN,标签 ready / mrv2\nmergeable,仅待审批;8-09 仍在更新"],
        ["规模", "28 文件 +1240/-110", "40 文件 +1507/-111"],
        ["范围", "变长投机的基础设施:\n图兼容、结构化输出、logprobs、\n混合 prefill/decode 批",
         "在此之上加入按置信度的预算调度\n(本项目对标的完整形态)"],
        ["验证模型", "Qwen3-8B-FP8 + DSpark", "DeepSeek-V4-Flash-DSpark,TP=4"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.09, 0.42, 0.49], row_h=0.62)

    tf = _txbox(s, Inches(0.55), Inches(4.55), Inches(12.3), Inches(2.6))
    _p(tf, "#47808 的机制(昇腾侧需要对齐的就是这些)", size=16, bold=True, first=True, before=0)
    _p(tf, "· 打分:Triton kernel 按“生存概率”(逐位 confidence 的累积乘积)对所有请求的 draft slot 统一排序,"
           "跨请求竞争 —— 自信请求的第 5 位可以压过犹豫请求的第 1 位,择优录取直到全局预算用尽;", size=13, before=8)
    _p(tf, "· 成本:启动时用 dummy step profile 出成本曲线;batch 级预算在 CPU 上用双缓冲的陈旧 confidence 计算"
           "(不引入同步),per-request 分配在 GPU 上用实时值;", size=13, before=4)
    _p(tf, "· 图:decode cudagraph 改为变长 —— 按 token 网格捕获、以 max_query_len 派发、每个捕获槽非空;"
           "要求后端上报 AttentionCGSupport.ALWAYS(PR 中把两个 MLA 后端由 UNIFORM_BATCH 升级为 ALWAYS);",
       size=13, before=4)
    _p(tf, "· 限制:仅支持带 confidence head 的 DSpark;不支持 LoRA / 流水并行 / 输出 logprobs。", size=13, before=4)
    _p(tf, "另有第二份参考实现:SGLang 的 DSpark 集成(2026-07-06)采用 packed ragged varlen、"
           "graph key 仅按总 token 数,padding 浪费更小 —— 可作为图侧的备选设计。", size=12.5, color=MUTE, before=8)

    # 4 我们的差异化 ──────────────────────────────────────────────────────────
    s = _slide(prs, "三、昇腾侧的三个真实缺口 —— 这才是本项目的工作量",
               "均已在源码层面核实,不是推测")
    _table(s, [
        ["#", "缺口", "证据", "工作量判断"],
        ["1", "confidence head 在昇腾侧\n尚未加载(短期缺项)",
         "vllm-ascend 的 deepseek_v4_draft.py 中\n_remap_dspark_name 对 confidence_head.*\n直接 return None;上游 main 原本同样如此,\n#47808 正是解除它的那个 PR",
         "小 —— 非难点。上游改动仅为解除一行\n丢弃逻辑 + 一个线性头。草稿侧我们已就位:\n训练已产出该头,且已按上游区分 bias"],
        ["2", "没有支持变长 decode 图的\n昇腾注意力后端",
         "上游依赖 AttentionCGSupport.ALWAYS;\n昇腾图模式当前按固定形状捕获",
         "★ 大 —— 本项目的主要技术风险\n与主要投入所在"],
        ["3", "预算调度 kernel 无昇腾实现",
         "上游为 Triton kernel;\n成本曲线依赖设备侧 profile",
         "中 —— 算法确定,需重写并\n在昇腾上重新标定成本曲线"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.04, 0.20, 0.42, 0.34], row_h=1.05)

    tf = _txbox(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(1.6))
    _p(tf, "我们的既有条件:DSV4-Flash DSpark 草稿已训练完成(5 epoch,接受长度达已发布草稿的 99.5%),"
           "训练→转换→部署→评测全链路可复现,并已具备无投机基线用于加速比测量。"
           "confidence head 训练时输入是 detach 的,可在冻结骨干上分钟级重新拟合。",
       size=13, color=MUTE, first=True, before=0)

    # 4b-1 技术补充(概念):为什么变长会掉出设备图 ─────────────────────────
    s = _slide(prs, "三(补充 1)、为什么「变长」会掉出设备图",
               "先讲清楚机制,再落到算子 —— 这是本项目最大技术风险的来源")
    tf = _txbox(s, Inches(0.55), Inches(1.56), Inches(12.3), Inches(1.0))
    _p(tf, "设备图(昇腾 ACL Graph,对应 GPU 侧的 CUDA Graph)是把一串算子按【固定的张量形状与地址】预先录制下来,"
           "推理时直接重放整段,省掉逐算子下发的开销。解码是「算子小、次数多」的场景,这部分开销占比很高,"
           "所以能否入图直接决定性能。", size=13.5, first=True, before=0)
    _p(tf, "而录制的前提,是每一步的形状都一样。", size=13.5, bold=True, color=ACCENT, before=5)

    for k, (px, pw, ttl, col, ks, note, tot) in enumerate([
        (0.55, 6.02, "固定 K = 5   →   形状固定,可入图", GOOD, [5, 5, 5, 5],
         "每个请求每步都验证 5 个 token。batch 内 query 长度一致、总 token 数每步相同,"
         "录制一次即可一直重放。", "总 token 数 = 4 × 5 = 20,每一步都是这个数"),
        (6.76, 6.02, "自适应 K_i   →   形状逐步变化,默认出图", RISK, [5, 2, 7, 1],
         "每个请求按自身置信度拿到不同的 K_i,下一步又是另一组。query 长度参差(ragged)、"
         "总 token 数每步都在变。", "总 token 数 = 5+2+7+1 = 15,下一步又是别的数")]):
        _shape(s, MSO_SHAPE.RECTANGLE, px, 2.70, pw, 2.86, None, col, 1.5)
        _tbox(s, px + 0.14, 2.78, pw - 0.28, 0.32, [(ttl, True, col, 12.5)])
        for r, kk in enumerate(ks):
            y = 3.22 + r * 0.38
            _tbox(s, px + 0.14, y - 0.04, 0.84, 0.28, [("请求 " + str(r + 1), False, MUTE, 9)])
            for j in range(kk):
                _shape(s, MSO_SHAPE.RECTANGLE, px + 1.02 + j * 0.32, y, 0.26, 0.22, col, BG, 0.5)
            _tbox(s, px + 1.02 + 7 * 0.32 + 0.12, y - 0.04, 0.9, 0.28, [("K=" + str(kk), False, col, 9)])
        _tbox(s, px + 0.16, 4.78, pw - 0.34, 0.36, [(note, False, INK, 10)])
        _fit(_shape(s, MSO_SHAPE.RECTANGLE, px + 0.14, 5.16, pw - 0.28, 0.30,
                    LGREY if k == 0 else RGBColor(0xFD, 0xEC, 0xEC)),
             [(tot, True, col, 10)], align=PP_ALIGN.CENTER, space=0)

    tf = _txbox(s, Inches(0.55), Inches(5.74), Inches(12.3), Inches(1.5))
    _p(tf, "⟹ 于是问题归结为一件很具体的事", size=15.5, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "图里那些【依赖序列长度】的参数,是在录制时就被固化成了常量,还是留成了运行时可以改内容的设备张量?"
           "前者一变就必须整张图重录,后者只改张量内容即可。答案因算子而异 —— 下一页逐个部件核对。",
       size=13.5, before=7)

    # 4b-2 技术补充(算子):逐部件核对 ─────────────────────────────────────
    s = _slide(prs, "三(补充 2)、落到算子:逐个部件核对",
               "四条结论均出自 vllm-ascend 现网代码,标注文件与行号,可自行核对")
    _table(s, [
        ["路径 / 部件", "机制", "代码出处(vllm-ascend)", "对变长 K 的影响"],
        ["通用路径\nFIA v2", "序列长度是\nhost 侧 Python list",
         "attention_v1.py:183\n  actual_seq_lengths_q: list[int]\nattention_v1.py:314  由\n  query_start_loc_cpu[1:].tolist() 构造",
         "值在捕获时被烘进图 —— K 一变图里的长度就是错的,\n必须为每个形状重新捕获。这就是「任意形状\n不可入图」的机械原因"],
        ["★ DSV4 路径\nDSA / SFA", "序列长度是\ndevice tensor",
         "device_op.py:559\n  actual_seq_lengths_query: torch.Tensor\ndsa_v1.py:1061  传入\n  query_start_loc[1:].clone()",
         "图捕获的是指针不是值 —— 改张量内容即可,结构上\n不需要重捕。★ 我们主攻的路径恰好在更有利的一侧"],
        ["workspace", "按 token 数做缓存键",
         "attention_v1.py:490/636/812/953\n  workspaces.get(num_tokens)\nattention_v1.py:774\n  num_tokens = actual_seq_lengths_q[-1]",
         "K 变 → num_tokens 变 → 缓存未命中、重新分配;\n不同 K 组合一多,图与 workspace 数量一起膨胀。\n★ 真正的成本在这里"],
        ["图内元数据\n更新", "机制已存在\n(固定 K 已在用)",
         "attention_v1.py:530-565\n  graph_task_update_begin/end 块内已有\n  _EXTRA_CTX.is_draft_model 与\n  attn_metadata[draft_step][key]",
         "固定 K 投机已经在图里跑,逐 draft step 更新图内\n元数据的机制与开销已经付过 —— 自适应不是从零\n造图内变长,而是让这套机制吃变化的长度"],
    ], Inches(0.55), Inches(1.56), Inches(12.3), col_w=[0.10, 0.15, 0.39, 0.36], row_h=0.80, font=10)
    _tbox(s, 0.55, 5.62, 12.3, 0.34,
          [("术语:FIA v2 = npu_fused_infer_attention_score_v2,昇腾融合注意力算子,通用模型走这条路径;"
            "DSA = DeepSeek Sparse Attention,SFA = 稀疏融合注意力,DeepSeek-V4-Flash 走这条;"
            "workspace = 算子运行所需的临时显存。", False, MUTE, 9)])

    tf = _txbox(s, Inches(0.55), Inches(6.02), Inches(12.3), Inches(1.3))
    _p(tf, "⟹ 准确的风险表述不是「昇腾不能变长入图」", size=15.5, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "而是:DSA 路径在设备侧持有长度、图内元数据更新机制已存在,两项结构性条件都已具备;"
           "未知量集中在 workspace 与图数量随 K 组合的膨胀,以及随之而来的显存占用与捕获耗时。", size=13, before=7)
    _p(tf, "这把阶段 2 的第一件事定死了:先量形状数量与 workspace 膨胀,再在两条路线里选 —— "
           "「对齐上游」按 token 总数分档捕获、运行时用批内最长请求的长度挑图并补齐(上游称后端上报 ALWAYS、"
           "按 max_query_len 派发);「SGLang 式」把各请求 token 紧挨着打包,图只按总 token 数分档,"
           "图少且无补齐浪费,但要改数据布局。先量再选,而不是先写实现。",
       size=13, bold=True, before=5)

    # 5 实施方案 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "四、实施方案与里程碑", "阶段 0 由“两周量化”缩为“一天判定”—— 上游数据已给出大部分答案")
    _table(s, [
        ["阶段", "周期", "主要工作", "交付物 / 决策门"],
        ["0\n判定", "1 天",
         "在昇腾现有固定 K 链路上,把加速比从 c=48 补测到\nc=128 / c=256,看是否出现上游那样的收益衰减",
         "一张并发-加速比曲线\nG1:是否值得继续"],
        ["1\n打通", "2–3 周",
         "vllm-ascend 加载 confidence head(移植上游改法);\n草稿侧供给并校验 confidence 数值合理性",
         "端到端可读出逐位置信度\nG2:置信度与实际接受率相关"],
        ["2\n图支持", "6–8 周",
         "★ 变长 decode 图:对齐 ALWAYS 语义或采用 SGLang\n的 packed varlen;预算 kernel 昇腾实现 + 成本曲线标定",
         "动态 K 在图模式下可运行\nG3:相对固定 K 不劣化"],
        ["3\n扩展", "3–4 周",
         "推广到通用模型路径;补齐限制项;文档与复现脚本",
         "可合入的实现 + 对比数据"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.07, 0.08, 0.47, 0.38], row_h=0.86)

    tf = _txbox(s, Inches(0.55), Inches(6.00), Inches(12.3), Inches(1.4))
    _p(tf, "总周期约 2.5–3.5 个月(v2 为 3–4 个月;算法设计部分由上游承担后收缩)",
       size=16, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "阶段 0 只要一天:我们已有无投机基线和固定 K 的完整评测链路,只需把并发点补齐。"
           "主要不确定性已从“收益幅度”转移到“昇腾图模式能否支持变长”—— 前者上游已用数据回答,后者是阶段 2 的核心。",
       size=13, before=8)

    # 6 风险 ──────────────────────────────────────────────────────────────────
    s = _slide(prs, "五、风险与依赖", "按“是否阻断”排序")
    _table(s, [
        ["风险", "影响", "应对"],
        ["★ 昇腾无法支持变长 decode 图\n(唯一的阻断性风险)",
         "阻断阶段 2,项目退化为\n“仅 eager 模式可用”",
         "阶段 2 先做可行性验证再投入;\n备选:SGLang 式 packed varlen,\ngraph key 只按总 token 数"],
        ["昇腾侧收益幅度不及 GPU",
         "项目价值下降,但不归零\n(同预算接受长度提升仍在)",
         "阶段 0 一天内给出判定;\nG1 不达标即停"],
        ["上游 #47808 设计仍在演进",
         "移植目标变动,返工",
         "跟随主线而非分叉;\n关闭的 #48692 已说明会向 #47808 收敛"],
        ["confidence head 质量不足",
         "调度信号噪声大,收益打折",
         "输入为 detach,可在冻结骨干上\n分钟级重训,代价极低"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.28, 0.30, 0.42], row_h=0.95)

    # 7 决策门 ────────────────────────────────────────────────────────────────
    s = _slide(prs, "六、决策门与验收标准", "每个门都是可量化的,不靠主观判断")
    _table(s, [
        ["门", "时点", "通过标准", "不通过怎么办"],
        ["G1", "第 1 天末",
         "在某个并发点上,固定 K 的加速比相对 c=48\n出现明显衰减(或已低于无投机)",
         "停止立项。结论本身有价值:\n说明昇腾当前运行点无需自适应"],
        ["G2", "第 3 周末",
         "端到端读出逐位置信度,且与实测逐位\n接受率单调相关",
         "回到草稿侧重新拟合 confidence head\n(代价分钟级),不影响主线"],
        ["G3", "第 11 周末",
         "动态 K 在图模式下运行,端到端不劣于\n同预算的固定 K",
         "保留 eager 路径成果并如实上报;\n将图支持作为独立议题提交上游"],
    ], Inches(0.55), Inches(1.62), Inches(12.3), col_w=[0.05, 0.11, 0.46, 0.38], row_h=0.98)

    tf = _txbox(s, Inches(0.55), Inches(5.20), Inches(12.3), Inches(2.0))
    _p(tf, "最终验收", size=16, bold=True, first=True, before=0)
    _p(tf, "· 昇腾侧可用实现,与固定 K、无投机三者的对比数据齐备,复现文档可被他人独立执行;", size=13.5, before=6)
    _p(tf, "· DSV4-Flash 在动态 K 下取得不劣于固定 K 的端到端表现,并在高并发点显示优势;", size=13.5, before=4)
    _p(tf, "· 通用模型路线给出明确结论:可行方案,或阻断原因与所需外部支持。", size=13.5, before=4)
    _p(tf, "建议:先做阶段 0。一天即可拿到 G1 数据,再决定是否投入后续 2.5–3.5 个月。",
       size=16, bold=True, color=ACCENT, before=12)

    prs.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="adaptive-spec-project-proposal-v3.pptx")
    build(ap.parse_args().out)
