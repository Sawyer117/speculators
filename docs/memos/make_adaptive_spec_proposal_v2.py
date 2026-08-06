#!/usr/bin/env python3
"""生成 v2 立项建议书 PPT:昇腾 NPU 自适应投机解码图模式支持。

与 v1 技术备忘录(make_adaptive_spec_memo.py)并存,不覆盖。
v1 = 调查结论;v2 = 按立项标准重写(范围/怎么做/多久/风险/决策门)。
用法:python3 make_adaptive_spec_proposal_v2.py --out <file>.pptx
"""
from __future__ import annotations

import argparse
import copy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1B, 0x22, 0x2E)
MUTE = RGBColor(0x5B, 0x66, 0x77)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
GOOD = RGBColor(0x0E, 0x7A, 0x4B)
WARN = RGBColor(0xB4, 0x53, 0x09)
RISK = RGBColor(0xA3, 0x1D, 0x2A)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF2, 0xF5, 0xF9)
CJK = "微软雅黑"


def _cjk(run, name=CJK):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        if latin is None:
            return
        ea = copy.deepcopy(latin)
        ea.tag = qn("a:ea")
        rPr.append(ea)
    ea.set("typeface", name)


def _txbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    return tf


def _p(tf, text, *, size=15, bold=False, color=INK, before=6, after=0,
       align=PP_ALIGN.LEFT, first=False, mono=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before, para.space_after = Pt(before), Pt(after)
    run = para.add_run()
    run.text = text
    run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color
    _cjk(run, "Consolas" if mono else CJK)
    return para


def _slide(prs, title, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.55), Inches(0.28), Inches(12.3), Inches(0.95))
    _p(tf, title, size=25, bold=True, first=True, before=0)
    if kicker:
        _p(tf, kicker, size=12, color=MUTE, before=3)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.20), Inches(1.3), Emu(24000))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background(); ln.shadow.inherit = False
    return s


def _table(slide, rows, x, y, w, col_w=None, font=11.5, row_h=0.32, colors=None, bolds=None):
    t = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, Inches(row_h * len(rows))).table
    if col_w:
        tot = sum(col_w)
        for i, f in enumerate(col_w):
            t.columns[i].width = Emu(int(w * f / tot))
    t.first_row = True
    for r, row in enumerate(rows):
        t.rows[r].height = Inches(row_h)
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if r == 0 else BG
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.text_frame.word_wrap = True
            for i, line in enumerate(str(val).split("\n")):
                para = cell.text_frame.paragraphs[0] if i == 0 else cell.text_frame.add_paragraph()
                para.space_before = para.space_after = Pt(0)
                run = para.add_run()
                run.text = line
                run.font.size = Pt(font)
                run.font.bold = (r == 0) or bool(bolds and (r, c) in bolds)
                run.font.color.rgb = colors[(r, c)] if (colors and (r, c) in colors) else INK
                _cjk(run)
    return t


def build(out: str) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1 封面 ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    tf = _txbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.0))
    _p(tf, "昇腾 NPU 自适应投机解码 · 图模式支持", size=36, bold=True, first=True, before=0)
    _p(tf, "立项建议书 v2", size=22, color=ACCENT, before=8)
    _p(tf, "覆盖通用模型(FIA 路径)与 DeepSeek-V4-Flash(DSA 路径)两条路线", size=17, color=MUTE, before=12)
    _p(tf, "技术判断已完成前置调查,结论出自 vllm-ascend 源码、vLLM PR #48692、昇腾算子文档;"
           "本文只讲清楚:问题是什么、怎么做、大概要多久", size=12, color=MUTE, before=20)

    # 2 立项理由 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "一、为什么立项", "投机解码的下一个增量不在训练侧,而在“每个请求验证多长”")
    _table(s, [
        ["", "现状", "问题", "立项要解决的"],
        ["算法侧", "全 batch 统一验证长度 K", "高置信请求本可多验,低置信请求\n白白浪费 target 计算", "按置信度给每请求分配 K_i"],
        ["工程侧", "固定 K 已入图,收益稳定", "K_i 每轮变化 → ragged batch,\n与设备图“形状固定”天然冲突", "让变长与图共存,且开销可控"],
        ["生态侧", "GPU 侧有原型(未合入)\nNPU 侧空白", "上游 PR 仅覆盖一种后端,\n且未证明加速", "NPU 侧首个可用实现,\n通用模型与 DSV4 都受益"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.10, 0.26, 0.34, 0.30], row_h=0.78)
    tf = _txbox(s, Inches(0.55), Inches(4.55), Inches(12.3), Inches(2.4))
    _p(tf, "收益随 target 变贵而放大 —— 这是本项目价值的核心判断", size=17, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "· GPU 侧原型在 8B 稠密模型上测不出加速(verify 本来就便宜),但同预算下接受长度从 3.445 提升到 3.707;", size=14, before=8)
    _p(tf, "· target 越大、verify 越贵,少验一个 token 省下的越多 —— 大 MoE(如 DSV4-Flash)是收益最大的场景;", size=14, before=4)
    _p(tf, "· 我们已有 DSpark 训练侧与 serve 侧的完整链路,是国内少数具备端到端验证条件的团队。", size=14, before=4)

    # 3 目标与范围 ────────────────────────────────────────────────────────────
    s = _slide(prs, "二、目标与范围", "明确不做什么,与明确做什么同样重要")
    _table(s, [
        ["", "做(In scope)", "不做(Out of scope)"],
        ["算法",
         "· confidence head 输出 → 每请求 K_i 分配\n· 总验证预算按 batch 规模确定",
         "· 不重新训练 confidence head\n· 不改动 draft 模型结构"],
        ["工程",
         "· K_i → 有效 token 的设备侧 compaction\n· workspace / 图键改为 bucket 化\n· 捕图期预分配与长期复用\n· 变长 metadata 与图 replay 打通",
         "· 不做 PD 分离、不改调度器\n· 不承诺 FIA v2 算子接口改动\n  (外部依赖,见风险页)"],
        ["范围",
         "· DSV4-Flash(DSA 路径)端到端\n· 通用稠密模型(FIA 路径)可行性与摊薄方案",
         "· 不覆盖 MTP / Eagle3 等其它投机算法\n· 不覆盖训练侧改动"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.09, 0.47, 0.44], row_h=1.05)
    tf = _txbox(s, Inches(0.55), Inches(5.5), Inches(12.3), Inches(1.4))
    _p(tf, "交付物:可用实现(vllm-ascend 分支)+ 四模式对比数据 + 上游 PR + 复现文档。",
       size=15, bold=True, first=True, before=0)

    # 4 可行性 ────────────────────────────────────────────────────────────────
    s = _slide(prs, "三、技术可行性 —— 前置调查已完成", "三条结论均有代码/文档出处,不是推测")
    _table(s, [
        ["", "调查结论", "对立项的意义"],
        ["设备侧长度",
         "DSA / lightning-indexer / SFA 的签名即 torch.Tensor,调用处传\nquery_start_loc[1:].clone(),全程不过 host;\nFIA v2 的 actual_seq_qlen 是 Host int 数组,全部调用点 .tolist()",
         "DSV4 前置条件已满足 → 可立即开工;\n通用模型需额外方案(见风险)"],
        ["增量代价",
         "固定-K 投机早已入图,replay 前逐层更新 seq lengths,\ngraph_task_update 开销已经在付;动态 K 的真正增量是\nworkspace 每轮 miss(缓存键=本轮总 token 数)+ 图键爆炸 + padding",
         "工作量收敛到 3 个明确点,\n不是“重写图路径”"],
        ["GPU 参考",
         "PR #48692 已关闭未合入;不做在线预算分配;仅 FLASH_ATTENTION;\n作者自述测不出加速",
         "不能直接移植;\n但变长与全图共存的设计可借鉴"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.13, 0.55, 0.32], row_h=1.02,
       colors={(1, 2): GOOD, (2, 2): GOOD, (3, 2): WARN})
    tf = _txbox(s, Inches(0.55), Inches(5.5), Inches(12.3), Inches(1.3))
    _p(tf, "结论:技术路径清晰,主要不确定性在“收益幅度”而非“能否实现” —— 因此第一阶段设为收益量化,而非直接开发。",
       size=15, bold=True, color=ACCENT, first=True, before=0)

    # 5 实施方案与里程碑 ──────────────────────────────────────────────────────
    s = _slide(prs, "四、实施方案与里程碑", "三阶段推进,每阶段结束设决策门")
    _table(s, [
        ["阶段", "周期", "主要工作", "交付物", "决策门"],
        ["阶段 0\n收益量化", "2 周",
         "· 统计同流量下静态 K 与自适应分配的\n  target 实际验证 token 数\n· 搭建四模式对比与 profiling 工具",
         "收益评估报告\n基线数据",
         "G1:token 减少量\n是否值得继续"],
        ["阶段 1\nDSV4 打通", "6–8 周",
         "· K_i 设备侧 compaction 全链路\n· workspace / 图键 bucket 化 + 捕图预分配\n· 端到端跑通、调优、四模式实测",
         "vllm-ascend 分支\n端到端性能数据",
         "G2:端到端时延\n是否优于静态 K 全图"],
        ["阶段 2\n通用扩展", "4–6 周",
         "· FIA 路径变长兼容标志对齐\n· host metadata 分桶摊薄方案\n· 推动算子侧接受设备张量(外部协同)",
         "通用模型支持\n上游 PR",
         "G3:是否推动\n算子接口改动"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.12, 0.08, 0.42, 0.20, 0.18], row_h=1.15)
    tf = _txbox(s, Inches(0.55), Inches(5.8), Inches(12.3), Inches(1.2))
    _p(tf, "总周期约 3–4 个月:阶段 0 的结论直接决定阶段 1 是否启动;阶段 2 可与阶段 1 部分并行,但算子改动依赖外部排期。",
       size=15, bold=True, color=ACCENT, first=True, before=0)
    _p(tf, "两周即可拿到第一个决策依据 —— 主要不确定性在“收益幅度”而非“能否实现”,所以先量化,不直接开发。",
       size=14, color=MUTE, before=8)

    # 7 风险与依赖 ────────────────────────────────────────────────────────────
    s = _slide(prs, "五、风险与依赖", "按“是否阻断”排序")
    _table(s, [
        ["风险 / 依赖", "影响", "应对"],
        ["收益幅度不足 —— 动态分配少算的 token\n不足以覆盖图外开销", "阻断,项目失去意义",
         "阶段 0 先量化;G1 不达标即停,\n只花两周,不进入开发"],
        ["FIA v2 不接受设备张量长度\n(通用模型路线)", "阻断通用模型,不影响 DSV4",
         "先交付 DSV4;并行以 bucket 摊薄 host\nmetadata;同时推动算子侧支持"],
        ["图键 / 捕图数量爆炸", "显存与启动时间上升",
         "按 (batch bucket × 总验证 token bucket)\n建键,限制桶数并实测 padding 浪费"],
        ["昇腾算子与 CANN 版本演进", "接口变动导致返工",
         "实现与算子调用之间保留派发层,\n跟随 vllm-ascend 主线节奏"],
        ["上游 GPU 侧方案变化\n(#48692 已关闭)", "参考失效,不阻断",
         "以我方实测为准,\n不以上游 PR 为前提"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.32, 0.24, 0.44], row_h=0.72,
       colors={(1, 1): RISK, (2, 1): WARN})

    # 8 验收标准 ──────────────────────────────────────────────────────────────
    s = _slide(prs, "六、决策门与验收标准", "每个门都是可量化的,不靠主观判断")
    _table(s, [
        ["门", "时点", "通过标准", "不通过则"],
        ["G1\n收益量化", "第 2 周",
         "同流量下自适应分配相对静态 K 的\ntarget 验证 token 数显著减少",
         "终止立项,\n输出评估报告"],
        ["G2\n端到端可用", "第 10 周",
         "动态 K 全图模式端到端 decode 时延\n优于静态 K 全图;图命中率与 padding\n浪费在可接受范围",
         "回退到静态 K,\n保留 profiling 成果"],
        ["G3\n通用扩展", "第 16 周",
         "通用稠密模型在 FIA 路径上可运行,\n或明确算子改动方案与排期",
         "范围收敛为\nDSV4 专用"],
    ], Inches(0.55), Inches(1.70), Inches(12.3), col_w=[0.11, 0.11, 0.48, 0.30], row_h=1.00)
    tf = _txbox(s, Inches(0.55), Inches(5.15), Inches(12.3), Inches(1.9))
    _p(tf, "最终验收", size=16, bold=True, first=True, before=0)
    _p(tf, "· 可用实现合入 vllm-ascend 分支,四模式对比数据齐备,复现文档可被他人独立执行;", size=14, before=6)
    _p(tf, "· 至少一个模型(DSV4-Flash)在动态 K 全图模式下取得端到端加速;", size=14, before=4)
    _p(tf, "· 通用模型路线给出明确结论:可行方案或阻断原因与所需外部支持。", size=14, before=4)
    _p(tf, "建议:先做阶段 0,用两周数据决定后面要不要继续。",
       size=16, bold=True, color=ACCENT, before=12)

    prs.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="adaptive-spec-project-proposal-v2.pptx")
    build(ap.parse_args().out)
