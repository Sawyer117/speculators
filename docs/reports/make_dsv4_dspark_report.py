#!/usr/bin/env python3
"""Generate the DSV4-DSpark-on-NPU project report deck (.pptx).

Self-contained (only python-pptx). Re-run to rebuild after editing content:
    python docs/reports/make_dsv4_dspark_report.py
-> writes docs/reports/dsv4_dspark_npu_report.pptx

Content is sourced from docs/deployment/ascend-npu-dsv4-dspark-ep-training.md (§1.5/§10 carry the
per-claim provenance); numbers are the measured run values as of 2026-07-14.
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INK = RGBColor(0x1B, 0x25, 0x38)      # dark slate
ACCENT = RGBColor(0x2E, 0x6C, 0xF6)   # blue
MUTE = RGBColor(0x5A, 0x66, 0x78)
GOOD = RGBColor(0x1B, 0x8A, 0x4E)
WARN = RGBColor(0xC0, 0x52, 0x00)
BG = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def _bg(slide, color=RGBColor(0xFF, 0xFF, 0xFF)):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Calibri"
    return r


def header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    tf = _tb(slide, 0.55, 0.14, 12.2, 0.95)
    if kicker:
        p = tf.paragraphs[0]; _run(p, kicker.upper(), 12, RGBColor(0x8F, 0xB6, 0xFF), bold=True)
    p = tf.add_paragraph() if kicker else tf.paragraphs[0]
    _run(p, title, 26, RGBColor(0xFF, 0xFF, 0xFF), bold=True)


def bullets(slide, items, top=1.5, left=0.6, width=12.1, size=16):
    tf = _tb(slide, left, top, width, 5.6)
    for it in items:
        if isinstance(it, tuple):
            text, lvl, *rest = it
            color = rest[0] if rest else INK
        else:
            text, lvl, color = it, 0, INK
        p = tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(6)
        bullet = "•  " if lvl == 0 else "–  "
        _run(p, bullet, size, ACCENT if lvl == 0 else MUTE, bold=True)
        _run(p, text, size if lvl == 0 else size - 1, color)
    return tf


def table(slide, headers, rows, top=1.55, left=0.6, width=12.1, colw=None, fs=13):
    nr, nc = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(0.4 * nr)).table
    if colw:
        for i, w in enumerate(colw):
            gt.columns[i].width = Inches(w)
    for j, htext in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = INK
        pr = c.text_frame.paragraphs[0]; _run(pr, htext, fs, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = BG if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)
            pr = c.text_frame.paragraphs[0]
            _run(pr, str(val), fs, INK)
    return gt


def footer(slide, n):
    tf = _tb(slide, 11.4, 7.02, 1.7, 0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"DSV4-DSpark · NPU · {n}", 9, MUTE)


slides = []

# 1 — title
s = prs.slides.add_slide(BLANK); _bg(s, INK)
tf = _tb(s, 0.9, 2.2, 11.5, 3)
p = tf.paragraphs[0]; _run(p, "DSV4-Flash DSpark Draft Training on Ascend NPU", 40, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
p = tf.add_paragraph(); _run(p, "Design decisions, implementation & measured results", 20, RGBColor(0x9F, 0xC0, 0xFF))
p = tf.add_paragraph(); p.space_before = Pt(18)
_run(p, "HF-native  ·  FSDP2 + Expert-Parallel(8) MoE  ·  grouped-GEMM  ·  online HS extraction", 15, RGBColor(0xC7, 0xD3, 0xE8))
p = tf.add_paragraph(); _run(p, "Sawyer117 / Austin Wen   ·   2026-07-14", 13, RGBColor(0x8F, 0xA3, 0xC0))
slides.append(s)

# 2 — goal
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Goal & scope", "why")
bullets(s, [
    ("Reproduce DeepSeek's DSpark speculative-decoding draft for DeepSeek-V4-Flash on Ascend NPU (A2, 8×64 GB).", 0),
    ("Target: match the released draft's acceptance length — AL 3.94 @ num_spec=5 (vllm-ascend PR #11196).", 0),
    ("Draft = 3 DSV4-native decoder layers (MLA + sink + 256-expert MoE + hyper-connections) + Markov/confidence heads.", 0),
    ("Trained via speculators (HF-native, FSDP2 — NOT megatron); greedy temp=0 end-to-end (gen/train/eval self-consistent).", 0),
    ("Two hard requirements that shaped everything: fit 256 experts on one A2 node, and get target hidden states out of vLLM.", 0),
])
footer(s, 2)

# 3 — pipeline
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "End-to-end pipeline", "system")
tf = _tb(s, 0.6, 1.4, 12.1, 0.9)
p = tf.paragraphs[0]
_run(p, "env build  →  rollout (greedy)  →  serve + HS dump  →  train (FSDP2+EP8)  →  eval", 18, ACCENT, bold=True)
bullets(s, [
    ("Coupling that makes it one chain:  rollout row index  =  HS file name  =  trainer sample key.", 0),
    ("Rollout: target generates prompt+response (greedy, max-tokens 3072) → Arrow dataset with loss_mask.", 0),
    ("Serve: DeepSeek-V4-Flash bf16, 2×A2 (TP8/DP2, EP OFF); HS_DUMP writes hs_<row>.safetensors to shared storage.", 0),
    ("Train: reads Arrow + rolling HS buffer (produce-one/consume-one); no offline HS store.", 0),
    ("Eval: serve verifier + trained draft, measure accept length vs the released-draft baseline.", 0),
], top=2.35)
footer(s, 3)

# 4 — architecture
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Architecture — verified from the weights, not analogy", "what")
table(s, ["", "Target (DeepSeek-V4-Flash, 43 L)", "Draft (DSpark, 3 mtp, γ=5)"],
      [["Attention", "MLA + per-head sink", "MLA + per-head sink"],
       ["Long-context", "SWA-128 + CSA(c4a) + HCA(c128a)", "SWA-128 only — NO DSA"],
       ["FFN", "256-expert MoE, top-6, +shared", "256-expert MoE, top-6, +shared"],
       ["Residual", "mHC (Manifold-Constrained Hyper-Conn.)", "mHC"],
       ["Heads", "—", "main_proj@mtp0; markov+confidence@mtp2"]],
      colw=[1.7, 5.4, 5.0])
tf = _tb(s, 0.6, 5.9, 12.1, 1.2)
p = tf.paragraphs[0]
_run(p, "Ground truth = released checkpoint weight keys + DSpark paper (arXiv:2607.05147). ", 13, MUTE)
_run(p, "The draft has NO compressor/indexer keys → it is dense sliding-window, not sparse. "
        "\"MLA+DSA+SWA\" for the draft is WRONG; correct = MLA + sink + SWA.", 13, WARN, italic=True)
footer(s, 4)

# 5 — HS extraction
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Design decision #1 — HS extraction: PR-based dumper, not the native connector", "decision")
bullets(s, [
    ("Native  extract_hidden_states  connector fails on DSV4 for three independent reasons:", 0),
    ("KV memory pathology — a CacheOnly cache co-sized with the real KV; DSV4's KV is hyper-compressed so it OOMs.", 1, WARN),
    ("vLLM version lock — crashes on our pinned 0.23.0; only fixed on 0.24 (which flips the v1→v2 runner).", 1, WARN),
    ("PD-disaggregation clash — kv_producer mode is incompatible with Ascend balance-scheduling.", 1, WARN),
    ("Chosen: Plan-B dumper — rides the DSV4 model's existing target-hidden buffer (vllm-ascend #11571).", 0, GOOD),
    ("A runner post-forward hook copies the buffer + post-norm hidden to CPU; writes the standard ArrowDataset format.", 1),
    ("Zero extra NPU KV, runs on the validated 0.23.0 serve, no PD-disagg. Validated end-to-end (3436 tok/s).", 1, GOOD),
])
footer(s, 5)

# 6 — EP refactor
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Design decision #2 — Expert-Parallel MoE (DTensor-native)", "decision")
bullets(s, [
    ("Problem: 256 experts × 3 layers. Per-shape MoE recompiles → ~160 s forward spikes; per-expert FSDP blocks fused grouped-GEMM.", 0),
    ("Rejected — plain-tensor experts via FSDP ignored_params: mixes plain + DTensor → clip / AdamW / checkpoint crash → patch whack-a-mole → not upstreamable.", 0, WARN),
    ("Chosen — GroupedExperts + Shard(0) DTensors (torchtitan-aligned):", 0, GOOD),
    ("stacked expert weights, each rank's slice a Shard(0) DTensor on the SAME mesh as FSDP → every param a uniform DTensor.", 1),
    ("optimizer / clip / DCP checkpoint need NO special-casing — all patches removed.", 1, GOOD),
    ("forward: all-to-all tokens → local grouped-GEMM (npu_grouped_matmul) → all-to-all back (autograd-aware).", 1),
    ("Effect: 256-expert step  ~160 s recompile spikes  →  steady ~370 ms forward.  Clean, checkpointable, upstreamable.", 0, GOOD),
])
footer(s, 6)

# 7 — block_size fix
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "The block_size off-by-one (caught mid-run)", "correctness")
bullets(s, [
    ("DSV4 DSpark drafts γ = 5 tokens per block (num_speculative_tokens = 5).", 0),
    ("BUT the speculators trainer's block_size is the block WIDTH incl the anchor (slot 0, loss-masked): drafts = block_size − 1.", 0),
    ("So training needs --block-size 6, not 5. Passing 5 silently drafts only 4 (logs showed position_1..4).", 0, WARN),
    ("Three-way verified: dflash/core.py:188 (\"first block position is the anchor, not emitted\"); vllm-ascend n_predict = dspark_block_size (no −1); Qwen3 cross-check (train BLOCK_SIZE=8 ⇔ release block_size=7).", 0),
    ("Fixed → --block-size 6; logs now show position_1..5 (5 draft positions). Window: train 133/5, infer 132/4.", 0, GOOD),
])
footer(s, 7)

# 8 — training results
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Training results — epoch 0 complete (block6 @ 196 anchors, EP8)", "results")
table(s, ["metric", "start", "end of epoch 0", "target"],
      [["loss", "2.26", "0.29", "↓"],
       ["accept_len", "1.0", "2.32 ↑ (max 2.71)", "3.94 (released)"],
       ["position_1..5 acc", "—", "0.65 / 0.49 / 0.37 / 0.31 / 0.27", "—"],
       ["confidence pred vs obs", "—", "0.20 vs 0.38", "calibrate"]],
      colw=[3.3, 1.9, 4.8, 2.1], top=1.55)
bullets(s, [
    ("Still learning at epoch end — recent-500 vs prior-500: loss↓, accept_len↑, full_acc↑ (not plateaued).", 0, GOOD),
    ("accept_len 2.32 after 1 epoch (of 20); climbing toward the released-draft 3.94 — needs more epochs.", 0),
    ("Memory 59 GB / 64 (92%). Two crash bugs fixed this run: --block-size 6 (off-by-one) and the val", 0),
    ("dataloader fork-at-epoch-boundary (num_workers=0); training now survives epoch boundaries + resumes.", 1),
    ("Converged per-benchmark numbers pending train→eval (needs the stacked→per-expert converter).", 0),
], top=4.15)
footer(s, 8)

# 9 — timing
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Throughput — steady is a lie; spikes eat 75% of wall-clock", "performance")
table(s, ["", "steady (spikes excluded)", "EFFECTIVE (real, incl spikes)"],
      [["step time", "~1.17 s", "~4.8 s  (4.1×)"],
       ["throughput", "~2.6k tok/s", "~630 tok/s"],
       ["1 epoch (4416 steps)", "(would be ~1.4 h)", "5.8 h"]],
      colw=[3.4, 4.3, 4.4], top=1.5)
tf = _tb(s, 0.6, 3.05, 12.1, 0.4)
p = tf.paragraphs[0]; _run(p, "Spike overhead ≈ 15,700 s / epoch (75% of wall-clock) breaks down as:", 15, INK, bold=True)
table(s, ["cause", "count × each", "cost", "fix"],
      [["MoE recompile (new token-count shapes)", "~623 × ~18 s", "~11.2k s (71%)", "fixed-shape MoE padding"],
       ["checkpoint save (EP-DCP gather)", "~10 × ~300 s", "~3.0k s (19%)", "raise --checkpoint-freq (→1/epoch)"],
       ["real HS starvation", "~2 × ~300 s", "~0.6k s (4%)", "serve throughput (minor)"]],
      colw=[5.0, 2.4, 2.4, 3.0], top=3.5, fs=12)
bullets(s, [
    ("Two levers → ~3× (5.8 h → ~2 h/epoch): fixed-shape MoE padding (recompiles) + save once/epoch (checkpoints).", 0, GOOD),
    ("Backward (~620 ms) dominates the steady step; HS fetch (~27 ms) is NOT the bottleneck (giant 'fetch' spikes were checkpoints).", 0),
], top=5.55)
footer(s, 9)

# 10 — attention op
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Draft attention op — SWA + non-causal + sink", "kernel")
table(s, ["op", "SWA", "non-causal", "sink", "backward", "use"],
      [["vllm-ascend SAS (inference)", "✓", "✓", "✓", "✗", "inference / gold reference"],
       ["SDPA / npu_fusion_attention", "✓", "✗", "✗", "✓", "— (missing sink + non-causal)"],
       ["einsum (current training path)", "✓", "✓", "✓", "✓ slow", "training now"],
       ["Triton kernel (built)", "✓", "✓", "✓", "✓", "training target — 5.96e-7 vs gold"]],
      colw=[4.2, 1.1, 1.7, 1.1, 1.5, 3.0], top=1.55, fs=12)
bullets(s, [
    ("No single fused op has all four — so training uses the (correct but slow) einsum; the Triton kernel fuses it.", 0),
    ("Triton kernel is done and validated (matches the gold reference at fp32 5.96e-7); needs integration into training.", 0, GOOD),
    ("Real inference dispatches to the SAS op (source verified non-causal; #11196 relaxed the causal asserts).", 0),
], top=4.5)
footer(s, 10)

# 11 — validation
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Correctness validation — \"it runs\" ≠ \"it's correct\"", "validation")
bullets(s, [
    ("Every module needs a correctness ORACLE, not just a smoke test. Status today:", 0),
    ("GREEN (gold-standard oracle): MoE grouped-GEMM & all-to-all parity, HF component parity, MLA ref, draft-attn vs gold, mHC Sinkhorn.", 0, GOOD),
    ("PENDING: assembled-draft numeric parity; structural weight-key parity vs the released draft; HS value-correctness.", 0, WARN),
    ("Silently-wrong gates to add: EP-invariance (EP=1 vs 8), gradcheck on hand-written backward, overfit-one-batch, serve↔train config guard.", 0, WARN),
    ("Green-check runner tiered CPU / single-NPU / needs-serve; the block-size bug is exactly the kind the config guard catches.", 0),
])
footer(s, 11)

# 12 — roadmap
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Remaining work / roadmap", "next")
bullets(s, [
    ("Throughput (biggest single win, ~3×): fixed-shape MoE padding (kills recompiles) + save once/epoch (done, one-liner).", 0, ACCENT),
    ("Critical path: stacked→per-expert weight converter (unblocks serve/eval) → full train→eval → fill converged numbers.", 0, ACCENT),
    ("Validation gates (cheap, CPU): structural parity · EP-invariance · overfit-one-batch · config guard.", 0),
    ("Perf follow-ups: argsort AiCore fix · gradient checkpointing (→512 anchors) · SWA Triton kernel into training.", 0),
    ("Upstream: carve the DTensor-native EP refactor into a clean GPU-safe PR.", 0),
    ("Deferred: native extract HS path (needs vLLM 0.24 + the DP0/head memory pathology solved).", 0, MUTE),
])
footer(s, 12)

# 13 — provenance
s = prs.slides.add_slide(BLANK); _bg(s); header(s, "Reproducibility & provenance", "sources")
bullets(s, [
    ("Every load-bearing claim is tied to a file:line / released config field / PR / paper — see the design doc §1.5 + §10.", 0),
    ("Design & results: docs/deployment/ascend-npu-dsv4-dspark-ep-training.md  (branch feat/dsv4-dspark).", 0),
    ("Pipeline index + per-stage archive: docs/deployment/ascend-npu-dsv4-dspark-pipeline.md.", 0),
    ("HS extraction deep-dive: ascend-npu-dsv4-hs-dumper-planB.md.", 0),
    ("Repos: Sawyer117/speculators (trainer) · Sawyer117/vllm-ascend @ feat/dsv4-hs-dumper (serve+HS) · non-causal-swa-triton-ascend (kernel).", 0),
    ("Log analysis: examples/ascend_npu_dflash/analyze_train_run.py (loss / accept / confidence / timing / spikes + plots).", 0),
])
footer(s, 13)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "dsv4_dspark_npu_report.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
